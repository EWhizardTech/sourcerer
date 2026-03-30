# app/services/parsing/strategies/ppt_parser.py

import base64
import logging
import re
from io import BytesIO
from typing import Any, Dict, List, Tuple

from app.services.parsing.base import BaseParser

logger = logging.getLogger(__name__)


class PPTParser(BaseParser):
    """
    PowerPoint parser (.pptx / .ppt).

    Primary:  python-pptx  – structured extraction per slide
    Fallback: zipfile + raw XML – plain-text extraction when pptx is corrupt
              or the file is an older binary .ppt that python-pptx cannot open.

    Extracts:
    - text  (all slides joined)
    - sections  (one per slide, heading = slide title)
    - tables  (one entry per table, Markdown-style grid)
    - images  (base64-encoded, keyed by slide + shape index)
    - lists  (bullet blocks per slide)
    """

    # ------------------------------------------------------------------ #
    # Public API
    # ------------------------------------------------------------------ #

    def parse(self, content: bytes, file_name: str) -> Dict[str, Any]:
        pptx_error: str | None = None

        # ---- primary path ------------------------------------------------
        try:
            result = self._parse_with_pptx(content, file_name)

            if result["text"].strip():
                logger.info(
                    "Parsed PPT %s via python-pptx: %d sections, %d tables, %d images",
                    file_name,
                    len(result["sections"]),
                    len(result["tables"]),
                    len(result["images"]),
                )
                return result

            pptx_error = "python-pptx returned empty text"
            logger.warning(
                "python-pptx returned empty text for %s, using XML fallback",
                file_name,
            )

        except Exception as exc:
            pptx_error = str(exc)
            logger.warning(
                "python-pptx parsing failed for %s, using XML fallback: %s",
                file_name,
                exc,
            )

        # ---- fallback path -----------------------------------------------
        try:
            text, sections = self._parse_with_xml_fallback(content)
            logger.info(
                "Parsed PPT %s via XML fallback: %d sections",
                file_name,
                len(sections),
            )
            return self._build_result(
                text=text,
                sections=sections,
                tables=[],
                lists=[],
                images=[],
                metadata={
                    "parser": "ppt",
                    "file_name": file_name,
                    "extraction_backend": "xml_fallback",
                    "pptx_error": pptx_error,
                },
            )

        except Exception as fallback_exc:
            logger.exception(
                "PPT parsing failed for %s in both primary and fallback: %s",
                file_name,
                fallback_exc,
            )
            return self._build_result(
                text="",
                sections=[],
                tables=[],
                lists=[],
                images=[],
                metadata={
                    "parser": "ppt",
                    "file_name": file_name,
                    "error": str(fallback_exc),
                    "pptx_error": pptx_error,
                },
            )

    # ------------------------------------------------------------------ #
    # Primary: python-pptx
    # ------------------------------------------------------------------ #

    def _parse_with_pptx(
        self, content: bytes, file_name: str
    ) -> Dict[str, Any]:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(BytesIO(content))

        all_text_parts: List[str] = []
        sections: List[Dict[str, Any]] = []
        tables: List[Dict[str, Any]] = []
        lists_: List[Dict[str, Any]] = []
        images: List[Dict[str, Any]] = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_title = self._extract_slide_title(slide)
            slide_texts: List[str] = []
            slide_bullets: List[str] = []

            for shape_idx, shape in enumerate(slide.shapes):

                # ---- TABLE -----------------------------------------------
                if shape.has_table:
                    md_table = self._table_to_markdown(shape.table)
                    tables.append(
                        {
                            "type": "table",
                            "content": md_table,
                            "slide_number": slide_idx,
                        }
                    )
                    slide_texts.append(md_table)
                    continue

                # ---- TEXT FRAME ------------------------------------------
                if shape.has_text_frame:
                    bullets, plain = self._extract_text_frame(shape.text_frame)

                    if plain.strip():
                        slide_texts.append(plain)

                    # Skip the title shape itself from bullet list
                    if bullets and shape != self._get_title_shape(slide):
                        slide_bullets.extend(bullets)
                    continue

                # ---- IMAGE -----------------------------------------------
                if (
                    shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                    or hasattr(shape, "image")
                ):
                    try:
                        image_entry = self._extract_image(
                            shape, slide_idx, shape_idx
                        )
                        if image_entry:
                            images.append(image_entry)
                    except Exception as img_exc:
                        logger.debug(
                            "Could not extract image slide=%d shape=%d: %s",
                            slide_idx,
                            shape_idx,
                            img_exc,
                        )

            # ---- slide notes ---------------------------------------------
            notes_text = self._extract_notes(slide)
            if notes_text:
                slide_texts.append(f"[Notes] {notes_text}")

            # ---- flush bullets -------------------------------------------
            if slide_bullets:
                lists_.append(
                    {
                        "type": "list",
                        "content": "\n".join(slide_bullets),
                        "slide_number": slide_idx,
                    }
                )

            # ---- build section for this slide ----------------------------
            section_content = "\n".join(slide_texts).strip()
            if slide_title or section_content:
                sections.append(
                    {
                        "type": "section",
                        "heading": slide_title,
                        "content": section_content,
                        "slide_number": slide_idx,
                    }
                )
                all_text_parts.append(
                    "\n".join(filter(None, [slide_title, section_content]))
                )

        text = "\n\n".join(all_text_parts).strip()

        return self._build_result(
            text=text,
            sections=sections,
            tables=tables,
            lists=lists_,
            images=images,
            metadata={
                "parser": "ppt",
                "file_name": file_name,
                "extraction_backend": "python-pptx",
                "total_slides": len(prs.slides),
            },
        )

    # ------------------------------------------------------------------ #
    # Fallback: zipfile + raw XML text extraction
    # ------------------------------------------------------------------ #

    def _parse_with_xml_fallback(
        self, content: bytes
    ) -> Tuple[str, List[Dict[str, Any]]]:
        """
        .pptx files are ZIP archives. Each slide is at ppt/slides/slideN.xml.
        We strip XML tags and collect readable text per slide.
        Works even when python-pptx cannot open the file.
        """
        import zipfile

        sections: List[Dict[str, Any]] = []
        all_texts: List[str] = []

        with zipfile.ZipFile(BytesIO(content)) as zf:
            slide_files = sorted(
                [n for n in zf.namelist() if re.match(r"ppt/slides/slide\d+\.xml", n)],
                key=lambda n: int(re.search(r"\d+", n).group()),
            )

            for slide_idx, slide_file in enumerate(slide_files, start=1):
                xml_bytes = zf.read(slide_file)
                slide_text = self._strip_xml_tags(xml_bytes.decode("utf-8", errors="ignore"))

                if slide_text.strip():
                    sections.append(
                        {
                            "type": "section",
                            "heading": f"Slide {slide_idx}",
                            "content": slide_text.strip(),
                            "slide_number": slide_idx,
                        }
                    )
                    all_texts.append(slide_text.strip())

        return "\n\n".join(all_texts), sections

    # ------------------------------------------------------------------ #
    # Helpers
    # ------------------------------------------------------------------ #

    @staticmethod
    def _strip_xml_tags(xml: str) -> str:
        """Remove XML markup and collapse whitespace."""
        text = re.sub(r"<[^>]+>", " ", xml)
        text = re.sub(r"\s+", " ", text)
        return text.strip()

    @staticmethod
    def _get_title_shape(slide):
        """Returns the title placeholder or None."""
        try:
            return slide.shapes.title
        except Exception:
            return None

    @classmethod
    def _extract_slide_title(cls, slide) -> str | None:
        shape = cls._get_title_shape(slide)
        if shape and shape.has_text_frame:
            title = shape.text_frame.text.strip()
            return title or None
        return None

    @staticmethod
    def _extract_text_frame(tf) -> Tuple[List[str], str]:
        """
        Returns (bullet_lines, plain_text).

        bullet_lines preserves '- ' prefix for list items.
        plain_text is the raw joined text.
        """
        bullets: List[str] = []
        lines: List[str] = []

        for para in tf.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            lines.append(text)

            # Treat indented paragraphs or paragraphs with bullet chars as bullets
            if para.level > 0 or text.startswith(("•", "-", "–", "●", "*")):
                # Normalise bullet character
                cleaned = re.sub(r"^[•\-–●*]\s*", "", text)
                bullets.append(f"- {cleaned}")
            else:
                bullets.append(text)

        return bullets, "\n".join(lines)

    @staticmethod
    def _extract_notes(slide) -> str:
        """Extracts speaker notes text."""
        try:
            notes_slide = slide.notes_slide
            if notes_slide and notes_slide.notes_text_frame:
                return notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass
        return ""

    @staticmethod
    def _table_to_markdown(table) -> str:
        """Converts a pptx Table object to a Markdown-style grid string."""
        rows: List[List[str]] = []

        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            rows.append(cells)

        if not rows:
            return ""

        col_widths = [
            max(len(rows[r][c]) for r in range(len(rows)))
            for c in range(len(rows[0]))
        ]

        def fmt_row(cells: List[str]) -> str:
            return "| " + " | ".join(c.ljust(w) for c, w in zip(cells, col_widths)) + " |"

        separator = "|-" + "-|-".join("-" * w for w in col_widths) + "-|"

        lines = [fmt_row(rows[0]), separator]
        for row in rows[1:]:
            lines.append(fmt_row(row))

        return "\n".join(lines)

    @staticmethod
    def _extract_image(shape, slide_idx: int, shape_idx: int) -> Dict[str, Any] | None:
        """Extracts image bytes and returns a normalised image entry."""
        image = shape.image
        image_bytes_b64 = base64.b64encode(image.blob).decode("utf-8")
        ext = image.ext or "png"

        return {
            "image_id": f"slide{slide_idx}_shape{shape_idx}.{ext}",
            "image_bytes": image_bytes_b64,
            "page_number": slide_idx,       # "page" = slide for consistency
            "content_type": image.content_type,
        }

    @staticmethod
    def _build_result(
        text: str,
        sections: List[Dict[str, Any]],
        tables: List[Dict[str, Any]],
        lists: List[Dict[str, Any]],
        images: List[Dict[str, Any]],
        metadata: Dict[str, Any],
    ) -> Dict[str, Any]:
        return {
            "text": text,
            "sections": sections,
            "tables": tables,
            "lists": lists,
            "images": images,
            "external_content": [],
            "metadata": metadata,
        }